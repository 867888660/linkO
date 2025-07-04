import os
import re
import logging
import tempfile
import shutil
from pptx import Presentation
from copy import deepcopy
import win32com.client
import pythoncom

# 打印版本号确认加载成功
print(pythoncom.__file__)

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Define the number of outputs and inputs
OutPutNum = 1
InPutNum = 4

# Initialize Outputs and Inputs arrays and assign names directly
Outputs = [{'Num': None, 'Context': None, 'Boolean': False, 'Kind': 'String', 'Id': f'Output{i + 1}', 'name': 'WebOutput', 'Link': 0} for i in range(OutPutNum)]
Inputs = [{'Num': None, 'Context': None, 'Boolean': False, 'Kind': 'String', 'Id': f'Input{i + 1}', 'Isnecessary': True, 'name': 'WebInput', 'Link': 0, 'IsLabel': False} for i in range(InPutNum)]

NodeKind = 'Normal'
InputIsAdd = False
OutputIsAdd = False
Lable = [{'Id': 'Label1', 'Kind': 'None'}]
FunctionIntroduction='组件功能（简述代码整体功能）\\n这是一个PPT模板处理节点，主要用于批量修改PPT文件中的文本内容。通过解析替换规则，在PPT的幻灯片、文本框和表格中查找指定关键字并替换为新内容，同时保持原有的文本格式（字体、大小、颜色、加粗等）。支持合并模式，可将新生成的PPT内容添加到现有文件末尾。\n\n代码功能摘要（概括核心算法或主要处理步骤）\\n核心算法包括：1）解析替换规则字符串，提取关键字和替换文本的映射关系；2）遍历PPT中的所有幻灯片、形状、文本框和表格单元格；3）对每个文本段落进行关键字匹配和替换；4）使用深度格式复制技术保留原始文本的字体属性（包括复杂的颜色格式）；5）重构文本runs以维持原有格式；6）支持文件合并功能，使用COM接口处理PPT文件合并。\n\n参数\\n```yaml\\ninputs:\\n  - name: PPt_FilePath\\n    type: file\\n    required: true\\n    description: 模板PPT文件的路径\\n  - name: Replace_Content\\n    type: string\\n    required: true\\n    description: 替换规则，格式如：\"要替换的文本\".replace\"替换后的文本\"，多条规则用##分隔\\n  - name: Save_Name\\n    type: string\\n    required: true\\n    description: 保存的文件名\\n  - name: Save_Path\\n    type: file\\n    required: true\\n    description: 保存的目标路径\\noutputs:\\n  - name: Result\\n    type: string\\n    description: 处理结果的详细日志信息，包含替换统计和错误信息\\n```\n\n运行逻辑（用 - 列表描写详细流程）\\n- 验证输入参数：检查模板PPT文件是否存在，确保保存路径目录存在\\n- 解析替换规则：按##分隔符拆分替换规则，使用正则表达式提取关键字和替换文本的映射关系\\n- 创建临时文件：复制模板PPT到临时文件进行安全处理\\n- 打开PPT文件：使用python-pptx库加载临时PPT文件\\n- 遍历幻灯片：对每张幻灯片进行文本处理\\n- 处理文本框：检查每个形状是否包含文本框，遍历段落和runs\\n- 执行文本替换：在段落文本中查找关键字并替换，统计替换次数\\n- 保留文本格式：使用copy_font_color函数精确复制字体属性，包括RGB颜色、主题颜色、方案颜色等\\n- 重构文本runs：根据替换结果重新构建文本runs，保持原有格式\\n- 处理表格：遍历表格中的每个单元格，对单元格文本执行相同的替换和格式保留操作\\n- 文件保存处理：检查目标文件是否存在，如存在则启用合并模式\\n- PPT合并：使用COM接口将新PPT内容添加到现有文件末尾\\n- 清理资源：删除临时文件，输出详细的处理日志和替换统计信息'

# Assign properties to Inputs and Outputs
Inputs[0]['name'] = 'PPt_FilePath'
Inputs[0]['Kind'] = 'String_FilePath'
Inputs[1]['name'] = 'Replace_Content'
Inputs[1]['Kind'] = 'String'
Inputs[2]['name'] = 'Save_Name'
Inputs[2]['Kind'] = 'String'
Inputs[3]['name'] = 'Save_Path'
Inputs[3]['Kind'] = 'String_FilePath'
Outputs[0]['name'] = 'Result'
Outputs[0]['Kind'] = 'String'


def copy_font_color(target_font, source_font, debug_log):
    """安全地复制字体颜色，处理不同类型的颜色对象"""
    if not hasattr(source_font, 'color') or source_font.color is None:
        debug_log.append("源字体没有颜色属性或颜色为None")
        return
    
    # 获取颜色类型名称
    color_type = source_font.color.__class__.__name__
    debug_log.append(f"颜色对象类型: {color_type}")
    
    try:
        # 处理RGB颜色
        if hasattr(source_font.color, 'rgb') and source_font.color.rgb:
            debug_log.append(f"复制RGB颜色: {source_font.color.rgb}")
            target_font.color.rgb = source_font.color.rgb
        # 处理主题颜色
        elif hasattr(source_font.color, 'theme_color') and source_font.color.theme_color is not None:
            debug_log.append(f"复制主题颜色: {source_font.color.theme_color}")
            target_font.color.theme_color = source_font.color.theme_color
        # 处理方案颜色
        elif hasattr(source_font.color, 'scheme_color') and source_font.color.scheme_color is not None:
            debug_log.append(f"复制方案颜色: {source_font.color.scheme_color}")
            target_font.color.scheme_color = source_font.color.scheme_color
        # 处理特殊情况：_color属性
        elif hasattr(source_font.color, '_color') and source_font.color._color:
            debug_log.append(f"从_color属性复制颜色: {source_font.color._color}")
            if isinstance(source_font.color._color, str):
                target_font.color.rgb = source_font.color._color
        else:
            debug_log.append("无法确定颜色类型，跳过颜色复制")
    except Exception as e:
        debug_log.append(f"复制颜色时出错: {str(e)}")
    
    # 打印颜色对象的所有属性
    debug_log.append(f"颜色对象属性: {dir(source_font.color)}")
    # 打印所有可能的值，帮助调试
    debug_log.append(f"_color属性: {getattr(source_font.color, '_color', 'Not available')}")
    debug_log.append(f"type属性: {getattr(source_font.color, 'type', 'Not available')}")
    debug_log.append(f"rgb属性: {getattr(source_font.color, 'rgb', 'Not available')}")
    debug_log.append(f"theme_color属性: {getattr(source_font.color, 'theme_color', 'Not available')}")
    debug_log.append(f"brightness属性: {getattr(source_font.color, 'brightness', 'Not available')}")
    
    try:
        if color_type == '_RGBColor' and hasattr(source_font.color, 'rgb'):
            debug_log.append(f"RGB值: {source_font.color.rgb}")
            target_font.color.rgb = source_font.color.rgb
        elif color_type == '_ThemeColor' and hasattr(source_font.color, 'theme_color'):
            debug_log.append(f"主题颜色: {source_font.color.theme_color}")
            target_font.color.theme_color = source_font.color.theme_color
        elif color_type == '_SchemeColor' and hasattr(source_font.color, 'scheme_color'):
            debug_log.append(f"方案颜色: {source_font.color.scheme_color}")
            target_font.color.scheme_color = source_font.color.scheme_color
        elif color_type == '_NoneColor':
            debug_log.append("检测到_NoneColor类型，跳过颜色复制")
            # 对于_NoneColor类型，不做任何操作
        elif color_type == 'ColorFormat':
            debug_log.append("检测到ColorFormat类型")
            # 处理ColorFormat类型
            if hasattr(source_font.color, 'rgb') and source_font.color.rgb is not None and source_font.color.rgb != '':
                debug_log.append(f"从ColorFormat复制RGB值: {source_font.color.rgb}")
                target_font.color.rgb = source_font.color.rgb
            elif hasattr(source_font.color, 'theme_color') and source_font.color.theme_color is not None:
                debug_log.append(f"从ColorFormat复制主题颜色: {source_font.color.theme_color}")
                target_font.color.theme_color = source_font.color.theme_color
            elif hasattr(source_font.color, '_color') and source_font.color._color is not None:
                debug_log.append(f"从ColorFormat的_color属性复制颜色: {source_font.color._color}")
                # 根据_color的类型决定如何复制
                if isinstance(source_font.color._color, str):
                    target_font.color.rgb = source_font.color._color
                else:
                    # 可能需要其他处理方式
                    debug_log.append(f"_color类型: {type(source_font.color._color)}")
            elif hasattr(source_font.color, 'type') and source_font.color.type == 1:  # 假设1代表白色
                debug_log.append("检测到可能的白色，设置为FFFFFF")
                target_font.color.rgb = "FFFFFF"
            else:
                debug_log.append("ColorFormat没有可用的RGB或theme_color值")
                # 尝试读取更多属性以便调试
                for attr_name in dir(source_font.color):
                    if not attr_name.startswith('__') and not callable(getattr(source_font.color, attr_name)):
                        try:
                            attr_value = getattr(source_font.color, attr_name)
                            debug_log.append(f"属性 {attr_name}: {attr_value}")
                        except:
                            debug_log.append(f"无法读取属性 {attr_name}")
        else:
            debug_log.append(f"未知颜色类型: {color_type}")
            # 可以根据需要添加更多颜色类型的处理
    except Exception as e:
        # 如果复制颜色时出错，打印错误但不中断处理
        debug_log.append(f"复制颜色时出错: {str(e)}")

def merge_presentations(existing_ppt_path, new_ppt_path, debug_log):
    """尝试使用更健壮的COM接口方法合并PPT"""
    import pythoncom
    import time
    import os
    import sys
    
    # 确保使用绝对路径
    existing_ppt_path = os.path.abspath(existing_ppt_path)
    new_ppt_path = os.path.abspath(new_ppt_path)
    
    # 检查文件是否存在
    if not os.path.exists(existing_ppt_path):
        debug_log.append(f"错误：现有PPT文件不存在: {existing_ppt_path}")
        return False
    
    if not os.path.exists(new_ppt_path):
        debug_log.append(f"错误：新PPT文件不存在: {new_ppt_path}")
        return False
    
    debug_log.append("文件路径验证通过")
    
    try:
        # 初始化COM，使用多线程模式
        pythoncom.CoInitializeEx(pythoncom.COINIT_APARTMENTTHREADED)
        debug_log.append("COM环境已初始化（多线程模式）")
        
        # 创建PowerPoint应用程序对象（尝试多次）
        import win32com.client
        ppt_app = None
        attempts = 0
        
        while attempts < 3 and ppt_app is None:
            try:
                debug_log.append(f"尝试创建PowerPoint应用实例 (尝试 #{attempts+1})")
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                break
            except Exception as e:
                debug_log.append(f"创建PowerPoint实例失败: {e}")
                attempts += 1
                time.sleep(1)  # 等待一秒再试
        
        if ppt_app is None:
            debug_log.append("无法创建PowerPoint应用程序实例")
            return False
        
        # 不要设置Visible属性，因为这可能导致错误
        debug_log.append("成功创建PowerPoint应用程序实例")
        
        # 打开文件
        try:
            debug_log.append(f"尝试打开现有PPT: {existing_ppt_path}")
            existing_ppt = ppt_app.Presentations.Open(existing_ppt_path, ReadOnly=False)
            debug_log.append("成功打开现有PPT")
            
            debug_log.append(f"尝试打开新PPT: {new_ppt_path}")
            new_ppt = ppt_app.Presentations.Open(new_ppt_path, ReadOnly=True)
            debug_log.append("成功打开新PPT")
        except Exception as e:
            debug_log.append(f"打开PPT文件失败: {e}")
            if ppt_app:
                ppt_app.Quit()
            return False
        
        # 处理幻灯片
        try:
            original_count = existing_ppt.Slides.Count
            new_slides_count = new_ppt.Slides.Count
            
            debug_log.append(f"现有PPT有 {original_count} 张幻灯片")
            debug_log.append(f"新PPT有 {new_slides_count} 张幻灯片")
            
            # 逐个复制幻灯片
            for i in range(1, new_slides_count + 1):
                debug_log.append(f"正在复制第 {i} 张幻灯片（共 {new_slides_count} 张）")
                
                # 使用InsertFromFile方法而不是复制粘贴
                try:
                    existing_ppt.Slides.InsertFromFile(new_ppt_path, original_count + i - 1, i, i)
                    debug_log.append(f"已插入第 {i} 张幻灯片")
                except Exception as e:
                    debug_log.append(f"插入幻灯片 {i} 失败: {e}")
                    # 继续尝试下一张幻灯片
            
            # 保存文件
            debug_log.append("尝试保存文件")
            existing_ppt.Save()
            debug_log.append("文件保存成功")
            
        except Exception as e:
            debug_log.append(f"处理幻灯片时出错: {e}")
            
        # 关闭文件和应用
        try:
            existing_ppt.Close()
            new_ppt.Close()
            debug_log.append("PPT文件已关闭")
        except Exception as e:
            debug_log.append(f"关闭PPT文件时出错: {e}")
        
        try:
            ppt_app.Quit()
            debug_log.append("PowerPoint应用已退出")
        except Exception as e:
            debug_log.append(f"退出PowerPoint应用时出错: {e}")
        
        return True
        
    except Exception as e:
        debug_log.append(f"合并过程中发生未预期错误: {e}")
        import traceback
        debug_log.append(traceback.format_exc())
        return False
        
    finally:
        # 反初始化COM
        try:
            pythoncom.CoUninitialize()
            debug_log.append("COM环境已反初始化")
        except:
            pass

def run_node(node):
    """节点主函数"""
    # 提取输入参数
    ppt_file_path = node['Inputs'][0]['Context']  # 模板PPT路径
    ppt_file_path = ppt_file_path
    replace_content = node['Inputs'][1]['Context']  # 替换规则
    save_name = node['Inputs'][2]['Context']  # 输出文件名
    save_path = node['Inputs'][3]['Context']  # 输出目录路径
    save_path = save_path
    
    # 用于收集详细的匹配过程信息
    debug_log = []
    
    try:
        # 确保保存路径存在
        if not os.path.exists(save_path):
            os.makedirs(save_path, exist_ok=True)
            debug_log.append(f"创建保存目录: {save_path}")
        
        # 检查模板文件是否存在
        if not os.path.exists(ppt_file_path):
            error_msg = f"错误: 模板文件不存在 - {ppt_file_path}"
            debug_log.append(error_msg)
            logging.error(error_msg)
            Outputs[0]['Context'] = "\n".join(debug_log)
            return Outputs
            
        debug_log.append(f"模板文件路径: {ppt_file_path}")
            
        # 确保文件扩展名正确
        if not save_name.lower().endswith(('.ppt', '.pptx')):
            save_name += '.pptx'
            debug_log.append(f"修正文件名: {save_name}")
            
        # 完整的输出文件路径
        output_file_path = os.path.join(save_path, save_name)
        debug_log.append(f"输出文件路径: {output_file_path}")
        
        # 解析替换内容
        replacements = {}
        debug_log.append("解析替换规则:")
        debug_log.append(f"原始替换规则文本: {replace_content}")
        
        if replace_content:
            # 按照 ## 分割不同的替换规则
            rules = replace_content.split('##')
            debug_log.append(f"分割后的规则数量: {len(rules)}")
            
            for i, rule in enumerate(rules):
                rule = rule.strip()
                if not rule:
                    debug_log.append(f"  规则 #{i+1}: 空规则，已跳过")
                    continue
                
                debug_log.append(f"  规则 #{i+1}: '{rule}'")
                    
                # 解析替换规则
                match = re.search(r'"([^"]+)".replace"([^"]*)"', rule)
                if match:
                    keyword, replacement = match.groups()
                    replacements[keyword] = replacement
                    debug_log.append(f"    解析成功: 将 '{keyword}' 替换为 '{replacement}'")
                else:
                    debug_log.append(f"    解析失败: 规则格式不符合要求")
        
        debug_log.append(f"共解析出 {len(replacements)} 条有效替换规则")
        
        # 创建临时文件以处理PPT
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_path = temp_file.name
        
        debug_log.append(f"创建临时文件: {temp_path}")
        
        # 首先复制模板文件到临时文件
        shutil.copy2(ppt_file_path, temp_path)
        debug_log.append("已将模板复制到临时文件")
        
        # 打开临时PPT文件进行编辑
        try:
            prs = Presentation(temp_path)
            debug_log.append(f"成功打开PPT文件，共 {len(prs.slides)} 张幻灯片")
        except Exception as e:
            debug_log.append(f"打开PPT文件失败: {str(e)}")
            Outputs[0]['Context'] = "\n".join(debug_log)
            return Outputs
        
        # 记录替换统计
        total_replacements = 0
        
        # 处理每个幻灯片
        for slide_idx, slide in enumerate(prs.slides):
            slide_replacements = 0
            debug_log.append(f"\n幻灯片 #{slide_idx+1}:")
            
            # 处理每个形状
            for shape_idx, shape in enumerate(slide.shapes):
                shape_replacements = 0
                
                # 检查形状是否有文本框
                if shape.has_text_frame:
                    debug_log.append(f"  形状 #{shape_idx+1}: 有文本框")
                    text_frame = shape.text_frame
                    
                    # 遍历段落
                    for para_idx, paragraph in enumerate(text_frame.paragraphs):
                        # 获取整个段落的文本（合并所有runs）
                        para_text = "".join([run.text for run in paragraph.runs])
                        debug_log.append(f"    段落 #{para_idx+1}: '{para_text[:50]}{'...' if len(para_text) > 50 else ''}'")
                        
                        if not para_text.strip():
                            debug_log.append(f"    段落为空，跳过")
                            continue
                            
                        # 检查整个段落是否包含任何关键字
                        para_replacements = 0
                        modified_para_text = para_text
                        
                        for keyword, replacement in replacements.items():
                            if keyword in para_text:
                                # 计算替换发生的次数
                                occurrences = para_text.count(keyword)
                                para_replacements += occurrences
                                
                                # 在整个段落文本中执行替换
                                modified_para_text = modified_para_text.replace(keyword, replacement)
                                
                                debug_log.append(f"      找到关键字 '{keyword}'，替换为 '{replacement}'，共 {occurrences} 处")
                            else:
                                debug_log.append(f"      未找到关键字 '{keyword}'")
                        
                        # 如果有替换发生，则重建段落
                        if para_replacements > 0:
                            debug_log.append(f"      整段替换前: '{para_text}'")
                            debug_log.append(f"      整段替换后: '{modified_para_text}'")
                            
                            # 保存原始段落的格式和属性
                            original_runs = paragraph.runs
                            
                            # 如果只有一个run或没有格式差异，使用简单替换方法
                            if len(original_runs) <= 1:
                                # 清空所有现有runs
                                for i in range(len(original_runs)-1, -1, -1):
                                    p = paragraph._p
                                    p.remove(original_runs[i]._r)
                                    
                                # 添加新的单一run
                                run = paragraph.add_run()
                                run.text = modified_para_text
                                
                                # 如果有原始格式，复制过来
                                if original_runs:
                                    # 复制字体格式
                                    if hasattr(original_runs[0], 'font'):
                                        run.font.name = original_runs[0].font.name
                                        run.font.size = original_runs[0].font.size
                                        run.font.bold = original_runs[0].font.bold
                                        run.font.italic = original_runs[0].font.italic
                                        run.font.underline = original_runs[0].font.underline
                                        if hasattr(original_runs[0], 'font'):
                                            copy_font_color(run.font, original_runs[0].font, debug_log)
                            else:
                                # 处理多个runs的情况，保留原始格式
                                # 清空所有现有runs
                                for i in range(len(original_runs)-1, -1, -1):
                                    p = paragraph._p
                                    p.remove(original_runs[i]._r)
                                
                                # 创建一个映射，记录原始文本中每个字符对应的run
                                char_to_run_map = []
                                current_pos = 0
                                
                                for run in original_runs:
                                    run_length = len(run.text)
                                    for i in range(run_length):
                                        char_to_run_map.append(run)
                                    current_pos += run_length
                                
                                # 如果原始文本长度与映射长度不匹配，使用第一个run的格式
                                if len(para_text) != len(char_to_run_map):
                                    new_run = paragraph.add_run()
                                    new_run.text = modified_para_text
                                    
                                    # 应用第一个run的格式
                                    if original_runs:
                                        first_run = original_runs[0]
                                        if hasattr(first_run, 'font'):
                                            new_run.font.name = first_run.font.name
                                            new_run.font.size = first_run.font.size
                                            new_run.font.bold = first_run.font.bold
                                            new_run.font.italic = first_run.font.italic
                                            new_run.font.underline = first_run.font.underline
                                            if hasattr(first_run.font, 'color') and hasattr(first_run.font.color, 'rgb'):
                                                new_run.font.color.rgb = first_run.font.color.rgb
                                else:
                                    # 分析替换前后的文本，确定每个替换的位置
                                    replacements_map = []
                                    for keyword, replacement in replacements.items():
                                        if keyword in para_text:
                                            start_idx = 0
                                            while True:
                                                idx = para_text.find(keyword, start_idx)
                                                if idx == -1:
                                                    break
                                                replacements_map.append((idx, idx + len(keyword), replacement))
                                                start_idx = idx + len(keyword)
                                    
                                    # 按照起始位置排序替换映射
                                    replacements_map.sort(key=lambda x: x[0])
                                    
                                    # 创建新的runs，保留原始格式
                                    current_pos = 0
                                    for start_idx, end_idx, replacement in replacements_map:
                                        # 添加替换前的文本
                                        if start_idx > current_pos:
                                            for i in range(current_pos, start_idx):
                                                run = paragraph.add_run()
                                                run.text = para_text[i]
                                                # 复制原始格式
                                                original_run = char_to_run_map[i]
                                                if hasattr(original_run, 'font'):
                                                    run.font.name = original_run.font.name
                                                    run.font.size = original_run.font.size
                                                    run.font.bold = original_run.font.bold
                                                    run.font.italic = original_run.font.italic
                                                    run.font.underline = original_run.font.underline
                                                    if hasattr(original_run.font, 'color') and hasattr(original_run.font.color, 'rgb'):
                                                        run.font.color.rgb = original_run.font.color.rgb
                                        
                                        # 添加替换文本，使用原始关键字位置的第一个字符的格式
                                        run = paragraph.add_run()
                                        run.text = replacement
                                        # 使用关键字第一个字符的格式
                                        original_run = char_to_run_map[start_idx]
                                        if hasattr(original_run, 'font'):
                                            run.font.name = original_run.font.name
                                            run.font.size = original_run.font.size
                                            run.font.bold = original_run.font.bold
                                            run.font.italic = original_run.font.italic
                                            run.font.underline = original_run.font.underline
                                            if hasattr(original_run.font, 'color') and hasattr(original_run.font.color, 'rgb'):
                                                run.font.color.rgb = original_run.font.color.rgb
                                        
                                        current_pos = end_idx
                                    
                                    # 添加剩余文本
                                    if current_pos < len(para_text):
                                        for i in range(current_pos, len(para_text)):
                                            run = paragraph.add_run()
                                            run.text = para_text[i]
                                            # 复制原始格式
                                            original_run = char_to_run_map[i]
                                            if hasattr(original_run, 'font'):
                                                run.font.name = original_run.font.name
                                                run.font.size = original_run.font.size
                                                run.font.bold = original_run.font.bold
                                                run.font.italic = original_run.font.italic
                                                run.font.underline = original_run.font.underline
                                                if hasattr(original_run.font, 'color') and hasattr(original_run.font.color, 'rgb'):
                                                    run.font.color.rgb = original_run.font.color.rgb
                            debug_log.append(f"      整段替换前: '{para_text}'")
                            debug_log.append(f"      整段替换后: '{modified_para_text}'")
                            
                            # 保存原始段落的格式和属性
                            original_runs = paragraph.runs
                            run_formats = []
                            
                            # 如果只有一个run或没有格式差异，使用简单替换方法
                            if len(original_runs) <= 1:
                                # 清空所有现有runs
                                for i in range(len(original_runs)-1, -1, -1):
                                    p = paragraph._p
                                    p.remove(original_runs[i]._r)
                                    
                                # 添加新的单一run
                                run = paragraph.add_run()
                                run.text = modified_para_text
                                
                                # 如果有原始格式，复制过来
                                if original_runs:
                                    # 复制字体格式
                                    if hasattr(original_runs[0], 'font'):
                                        run.font.name = original_runs[0].font.name
                                        run.font.size = original_runs[0].font.size
                                        run.font.bold = original_runs[0].font.bold
                                        run.font.italic = original_runs[0].font.italic
                                        run.font.underline = original_runs[0].font.underline
                                        if hasattr(original_runs[0], 'font'):
                                            copy_font_color(run.font, original_runs[0].font,debug_log)
                            else:
                                # 这种情况更复杂，需要尽可能保留原始格式
                                # 保存原始的run和它们的文本位置映射关系
                                original_text_positions = []
                                current_pos = 0
                                
                                for run in original_runs:
                                    run_length = len(run.text)
                                    original_text_positions.append((current_pos, current_pos + run_length, run))
                                    current_pos += run_length
                                
                                # 清空所有现有runs
                                for i in range(len(original_runs)-1, -1, -1):
                                    p = paragraph._p
                                    p.remove(original_runs[i]._r)
                                
                                # 尝试映射替换后的文本到原始格式
                                # 简单起见，我们将使用一个单一的run
                                new_run = paragraph.add_run()
                                new_run.text = modified_para_text
                                
                                # 应用第一个run的格式作为基本格式
                                if original_runs:
                                    first_run = original_runs[0]
                                    # 复制字体格式
                                    if hasattr(first_run, 'font'):
                                        new_run.font.name = first_run.font.name
                                        new_run.font.size = first_run.font.size
                                        new_run.font.bold = first_run.font.bold
                                        new_run.font.italic = first_run.font.italic
                                        new_run.font.underline = first_run.font.underline
                                        # Use this conditional check:
                                        if hasattr(first_run.font, 'color'):
                                            if hasattr(first_run.font.color, 'rgb'):
                                                new_run.font.color.rgb = first_run.font.color.rgb
                                            elif hasattr(first_run.font.color, 'theme_color'):
                                                # For scheme colors, copy the theme color
                                                new_run.font.color.theme_color = first_run.font.color.theme_color
                                            # You can add more conditions for other color types if needed
                            
                            debug_log.append(f"    段落重构完成，替换了 {para_replacements} 处")
                            shape_replacements += para_replacements
                        else:
                            debug_log.append(f"    段落未发生变化，跳过重构")
                
                # 检查表格
                if shape.has_table:
                    table = shape.table
                    debug_log.append(f"  形状 #{shape_idx+1}: 有表格 ({len(table.rows)}行x{len(table.columns)}列)")
                    
                    # 遍历每个单元格
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            # 遍历单元格中的每个段落
                            for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                # 获取整个段落的文本（合并所有runs）
                                para_text = "".join([run.text for run in paragraph.runs])
                                
                                if not para_text.strip():
                                    continue
                                    
                                debug_log.append(f"    单元格 [{row_idx+1},{col_idx+1}] 段落 #{para_idx+1}: '{para_text[:30]}{'...' if len(para_text) > 30 else ''}'")
                                
                                # 检查整个段落是否包含任何关键字
                                cell_replacements = 0
                                modified_para_text = para_text
                                
                                for keyword, replacement in replacements.items():
                                    if keyword in para_text:
                                        # 计算替换发生的次数
                                        occurrences = para_text.count(keyword)
                                        cell_replacements += occurrences
                                        
                                        # 在整个段落文本中执行替换
                                        modified_para_text = modified_para_text.replace(keyword, replacement)
                                        
                                        debug_log.append(f"      找到关键字 '{keyword}'，替换为 '{replacement}'，共 {occurrences} 处")
                                    else:
                                        debug_log.append(f"      未找到关键字 '{keyword}'")
                                
                                # 如果有替换发生，则重建段落
                                if cell_replacements > 0:
                                    debug_log.append(f"      单元格段落替换前: '{para_text}'")
                                    debug_log.append(f"      单元格段落替换后: '{modified_para_text}'")
                                    
                                    # 保存原始段落的格式和属性
                                    original_runs = paragraph.runs
                                    
                                    # 如果只有一个run，使用简单替换方法
                                    if len(original_runs) <= 1:
                                        # 清空所有现有runs
                                        for i in range(len(original_runs)-1, -1, -1):
                                            p = paragraph._p
                                            p.remove(original_runs[i]._r)
                                            
                                        # 添加新的单一run
                                        run = paragraph.add_run()
                                        run.text = modified_para_text
                                        
                                        # 如果有原始格式，复制过来
                                        if original_runs:
                                            # 复制字体格式
                                            if hasattr(original_runs[0], 'font'):
                                                run.font.name = original_runs[0].font.name
                                                run.font.size = original_runs[0].font.size
                                                run.font.bold = original_runs[0].font.bold
                                                run.font.italic = original_runs[0].font.italic
                                                run.font.underline = original_runs[0].font.underline
                                                if hasattr(original_runs[0], 'font') and hasattr(original_runs[0].font, 'color'):
                                                    source_color = original_runs[0].font.color
                                                    color_type = source_color.__class__.__name__
                                                    
                                                    try:
                                                        if color_type == '_RGBColor' and hasattr(source_color, 'rgb'):
                                                            run.font.color.rgb = source_color.rgb
                                                        elif color_type == '_ThemeColor' and hasattr(source_color, 'theme_color'):
                                                            run.font.color.theme_color = source_color.theme_color
                                                        elif color_type == '_SchemeColor' and hasattr(source_color, 'scheme_color'):
                                                            run.font.color.scheme_color = source_color.scheme_color
                                                        # 如果是 _NoneColor 或其他类型，就不复制颜色属性
                                                    except Exception:
                                                        # 捕获任何可能的异常，确保程序不会因此中断
                                                        pass

                                    else:
                                        # 处理多个runs的情况，保留原始格式
                                        # 清空所有现有runs
                                        for i in range(len(original_runs)-1, -1, -1):
                                            p = paragraph._p
                                            p.remove(original_runs[i]._r)
                                        
                                        # 创建一个映射，记录原始文本中每个字符对应的run
                                        char_to_run_map = []
                                        current_pos = 0
                                        
                                        for run in original_runs:
                                            run_length = len(run.text)
                                            for i in range(run_length):
                                                char_to_run_map.append(run)
                                            current_pos += run_length
                                        
                                        # 如果原始文本长度与映射长度不匹配，使用第一个run的格式
                                        if len(para_text) != len(char_to_run_map):
                                            new_run = paragraph.add_run()
                                            new_run.text = modified_para_text
                                            
                                            # 应用第一个run的格式
                                            if original_runs:
                                                first_run = original_runs[0]
                                                if hasattr(first_run, 'font'):
                                                    new_run.font.name = first_run.font.name
                                                    new_run.font.size = first_run.font.size
                                                    new_run.font.bold = first_run.font.bold
                                                    new_run.font.italic = first_run.font.italic
                                                    new_run.font.underline = first_run.font.underline
                                                    if hasattr(first_run.font, 'color') and hasattr(first_run.font.color, 'rgb'):
                                                        new_run.font.color.rgb = first_run.font.color.rgb
                                        else:
                                            # 分析替换前后的文本，确定每个替换的位置
                                            replacements_map = []
                                            for keyword, replacement in replacements.items():
                                                if keyword in para_text:
                                                    start_idx = 0
                                                    while True:
                                                        idx = para_text.find(keyword, start_idx)
                                                        if idx == -1:
                                                            break
                                                        replacements_map.append((idx, idx + len(keyword), replacement))
                                                        start_idx = idx + len(keyword)
                                            
                                            # 按照起始位置排序替换映射
                                            replacements_map.sort(key=lambda x: x[0])
                                            
                                            # 创建新的runs，保留原始格式
                                            current_pos = 0
                                            for start_idx, end_idx, replacement in replacements_map:
                                                # 添加替换前的文本
                                                if start_idx > current_pos:
                                                    for i in range(current_pos, start_idx):
                                                        run = paragraph.add_run()
                                                        run.text = para_text[i]
                                                        # 复制原始格式
                                                        original_run = char_to_run_map[i]
                                                        if hasattr(original_run, 'font'):
                                                            run.font.name = original_run.font.name
                                                            run.font.size = original_run.font.size
                                                            run.font.bold = original_run.font.bold
                                                            run.font.italic = original_run.font.italic
                                                            run.font.underline = original_run.font.underline
                                                            if hasattr(original_run.font, 'color') and hasattr(original_run.font.color, 'rgb'):
                                                                run.font.color.rgb = original_run.font.color.rgb
                                                
                                                # 添加替换文本，使用原始关键字位置的第一个字符的格式
                                                run = paragraph.add_run()
                                                run.text = replacement
                                                # 使用关键字第一个字符的格式
                                                original_run = char_to_run_map[start_idx]
                                                if hasattr(original_run, 'font'):
                                                    run.font.name = original_run.font.name
                                                    run.font.size = original_run.font.size
                                                    run.font.bold = original_run.font.bold
                                                    run.font.italic = original_run.font.italic
                                                    run.font.underline = original_run.font.underline
                                                    if hasattr(original_run.font, 'color') and hasattr(original_run.font.color, 'rgb'):
                                                        run.font.color.rgb = original_run.font.color.rgb
                                                
                                                current_pos = end_idx
                                            
                                            # 添加剩余文本
                                            if current_pos < len(para_text):
                                                for i in range(current_pos, len(para_text)):
                                                    run = paragraph.add_run()
                                                    run.text = para_text[i]
                                                    # 复制原始格式
                                                    original_run = char_to_run_map[i]
                                                    if hasattr(original_run, 'font'):
                                                        run.font.name = original_run.font.name
                                                        run.font.size = original_run.font.size
                                                        run.font.bold = original_run.font.bold
                                                        run.font.italic = original_run.font.italic
                                                        run.font.underline = original_run.font.underline
                                                        if hasattr(original_run.font, 'color') and hasattr(original_run.font.color, 'rgb'):
                                                            run.font.color.rgb = original_run.font.color.rgb
                                    debug_log.append(f"      单元格段落替换前: '{para_text}'")
                                    debug_log.append(f"      单元格段落替换后: '{modified_para_text}'")
                                    
                                    # 保存原始段落的格式和属性
                                    original_runs = paragraph.runs
                                    
                                    # 清空所有现有runs
                                    for i in range(len(original_runs)-1, -1, -1):
                                        p = paragraph._p
                                        p.remove(original_runs[i]._r)
                                        
                                    # 添加新的单一run
                                    run = paragraph.add_run()
                                    run.text = modified_para_text
                                    
                                    # 如果有原始格式，复制过来
                                    if original_runs:
                                        # 复制字体格式
                                        if hasattr(original_runs[0], 'font'):
                                            # 基本属性复制
                                            run.font.name = original_runs[0].font.name
                                            run.font.size = original_runs[0].font.size
                                            run.font.bold = original_runs[0].font.bold
                                            run.font.italic = original_runs[0].font.italic
                                            run.font.underline = original_runs[0].font.underline
                                            
                                            # 使用改进的颜色复制函数
                                            copy_font_color(run.font, original_runs[0].font, debug_log)
                                    
                                    debug_log.append(f"    单元格段落重构完成，替换了 {cell_replacements} 处")
                                    shape_replacements += cell_replacements
                                else:
                                    debug_log.append(f"    单元格段落未发生变化，跳过重构")
                
                if not shape.has_text_frame and not shape.has_table:
                    debug_log.append(f"  形状 #{shape_idx+1}: 无文本框或表格，跳过")
                
                slide_replacements += shape_replacements
                if shape_replacements > 0:
                    debug_log.append(f"  形状共进行了 {shape_replacements} 次替换")
            
            total_replacements += slide_replacements
            debug_log.append(f"幻灯片 #{slide_idx+1} 共进行了 {slide_replacements} 次替换")
        
        debug_log.append(f"\n所有幻灯片共进行了 {total_replacements} 次替换")
        
        # 检查目标文件是否已存在
        file_exists = os.path.exists(output_file_path)
        
        if file_exists:
            debug_log.append(f"发现目标文件已存在: {output_file_path}")
            debug_log.append("准备合并模式: 将新PPT内容添加到现有PPT末尾")
            
            # 创建临时文件用于保存新生成的PPT
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as merged_temp_file:
                merged_temp_path = merged_temp_file.name
            
            debug_log.append(f"创建合并临时文件: {merged_temp_path}")
            
            # 保存新生成的PPT到临时路径
            try:
                prs.save(merged_temp_path)
                debug_log.append(f"成功保存新PPT到临时文件: {merged_temp_path}")
            except Exception as e:
                debug_log.append(f"保存新PPT到临时文件失败: {str(e)}")
                Outputs[0]['Context'] = "\n".join(debug_log)
                return Outputs
            
            # 合并PPT文件
            if merge_presentations(output_file_path, merged_temp_path, debug_log):
                debug_log.append("成功将新PPT合并到现有PPT文件末尾")
            else:
                debug_log.append("合并PPT文件失败")
            
            # 清理临时文件
            try:
                os.unlink(merged_temp_path)
                debug_log.append("成功删除合并临时文件")
            except Exception as e:
                debug_log.append(f"删除合并临时文件失败: {str(e)}")
        else:
            # 保存修改后的PPT到目标路径
            try:
                prs.save(output_file_path)
                debug_log.append(f"成功保存修改后的PPT到: {output_file_path}")
            except Exception as e:
                debug_log.append(f"保存PPT文件失败: {str(e)}")
                Outputs[0]['Context'] = "\n".join(debug_log)
                return Outputs
        
        # 清理临时文件
        try:
            os.unlink(temp_path)
            debug_log.append("成功删除临时文件")
        except Exception as e:
            debug_log.append(f"删除临时文件失败: {str(e)}")
        
        if total_replacements > 0:
            debug_log.append(f"成功: 已生成修改后的PPT文件，共执行 {total_replacements} 次替换")
        else:
            debug_log.append("警告: 未执行任何替换操作，请检查关键字是否正确")
        
        Outputs[0]['Context'] = "\n".join(debug_log)
        
    except Exception as e:
        error_message = f"处理PPT模板时出错: {str(e)}"
        logging.error(error_message)
        debug_log.append(error_message)
        import traceback
        debug_log.append(traceback.format_exc())
        Outputs[0]['Context'] = "\n".join(debug_log)
    
    return Outputs

