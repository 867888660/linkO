import json
import http.client
import requests
import re
import copy
import chardet
import logging
import os
import math
import io
from pptx import Presentation
from PyPDF2 import PdfReader  # 使用PyPDF2替代fitz

# **Define the number of outputs and inputs**
OutPutNum = 4
InPutNum = 3
# **Define the number of outputs and inputs**

# **Initialize Outputs and Inputs arrays and assign names directly**
Outputs = [{'Num': None, 'Kind': None, 'Boolean': False, 'Id': f'Output{i + 1}', 'Context': None, 'name': f'OutPut{i + 1}', 'Link': 0, 'Description': ''} for i in range(OutPutNum)]
Inputs = [{'Num': None, 'Kind': None, 'Id': f'Input{i + 1}', 'Context': None, 'Isnecessary': True, 'name': f'Input{i + 1}', 'Link': 0, 'IsLabel': True} for i in range(InPutNum)]
NodeKind = 'ArrayTrigger'
Lable = [{'Id': 'Label1', 'Kind': 'None'}]
FunctionIntroduction='组件功能（简述代码整体功能）\\n这是一个多格式文件分割处理器，能够读取并分割不同类型的文件内容（支持txt、pdf、ppt格式），将大型文档按指定规则分割成小片段进行处理。\n\n代码功能摘要（概括核心算法或主要处理步骤）\\n核心算法包括：文件类型自动识别、多编码格式兼容读取、灵活的分割规则解析（支持自定义分隔符、页面分割、正则表达式模式）、内容分组处理、以及结构化输出生成。主要处理步骤为文件读取→内容分割→分组处理→结果输出。\n\n参数\\n```yaml\\ninputs:\\n  - name: file_path\\n    type: file\\n    required: true\\n    description: 待分割的文件路径，支持txt、pdf、ppt格式\\n  - name: divide_sign\\n    type: string\\n    required: false\\n    description: 分割符号，支持普通字符串、\"Page\"模式或\"字符1@<X>字符2\"自定义格式\\n    default: \"\\\\n\\\\n\"\\n  - name: Max_ArrayNum\\n    type: string\\n    required: false\\n    description: 每组包含的片段数量，用于将分割结果进行分组\\n    default: \"1\"\\noutputs:\\n  - name: context\\n    type: string\\n    description: 分割后的文本内容片段\\n  - name: num\\n    type: integer\\n    description: 当前片段在结果数组中的序号\\n  - name: Is_Last\\n    type: boolean\\n    description: 标识当前片段是否为最后一个片段\\n  - name: Total_Num\\n    type: integer\\n    description: 分割后的总片段数量\\n```\n\n运行逻辑（用 - 列表描写详细流程）\\n- 读取输入参数：文件路径、分割符号和分组大小\\n- 根据文件扩展名识别文件类型（.txt/.pdf/.ppt/.pptx）\\n- 针对不同文件类型采用相应的读取方式：PDF使用PdfReader逐页提取文本，PPT使用Presentation读取幻灯片内容，文本文件尝试多种编码格式读取\\n- 解析分割规则：若divide_sign为\"Page\"则按页面分割，若包含\"@<>\"格式则解析为自定义正则表达式模式，否则使用普通字符串分割\\n- 根据分割规则对文件内容进行分割：自定义模式使用正则表达式分割并保留分隔符，其他情况使用字符串split方法\\n- 过滤空白内容，移除只包含空格或换行的片段\\n- 根据Max_ArrayNum参数对分割结果进行分组，将多个片段合并为一个输出单元\\n- 为每个分割片段生成完整的输出信息：包含内容、序号、是否最后一个、总数量等属性\\n- 返回包含所有分割结果的结构化数组，每个元素包含四个输出字段的完整信息'

# **Assign properties to Inputs**
Outputs[0]['Kind'] = 'String'
Outputs[1]['Kind'] = 'Num'
Outputs[0]['name'] = 'context'
Outputs[1]['name'] = 'num'
Outputs[2]['name'] = 'Is_Last'
Outputs[2]['Kind'] = 'Boolean'
Outputs[3]['name'] = 'Total_Num'
Outputs[3]['Kind'] = 'Num'
for input in Inputs:
    input['Kind'] = 'String_FilePath'
Inputs[0]['name'] = 'file_path'
Inputs[1]['name'] = 'divide_sign'
Inputs[1]['Isnecessary'] = False
Inputs[1]['Kind'] = 'String'
Inputs[2]['name'] = 'Max_ArrayNum'
Inputs[2]['Kind'] = 'String'
Inputs[2]['Isnecessary'] = False
# **Assign properties to Inputs**

def parse_custom_pattern(divide_sign):
    # 检查是否包含自定义模式
    if divide_sign and '@<' in divide_sign and '>' in divide_sign:
        # 提取字符1和字符2
        parts = divide_sign.split('@<')
        char1 = parts[0]
        char2 = parts[1].split('>')[1]
        # 返回用于匹配的正则表达式模式
        return f"{char1}[^\n]*?{char2}"
    return None

def split_by_custom_pattern(content, pattern):
    # 使用正则表达式分割内容
    # 保留分隔符
    splits = re.split(f"({pattern})", content)
    # 过滤空白内容
    filtered_splits = []
    current_content = ""
    
    for i, split in enumerate(splits):
        if re.match(pattern, split):  # 如果是分隔符
            if current_content:  # 如果有累积的内容
                filtered_splits.append(current_content.strip())
            current_content = ""  # 重置累积内容
        else:
            current_content += split  # 累积内容
    
    # 添加最后一段内容
    if current_content:
        filtered_splits.append(current_content.strip())
    
    return [s for s in filtered_splits if s.strip()]

def extract_text_from_pdf(file_path):
    """从PDF提取文本，返回每页的文本内容列表"""
    try:
        # 使用PyPDF2打开PDF文件
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            pages = []
            
            # 逐页提取文本
            for page in reader.pages:
                text = page.extract_text()
                pages.append(text)
                
            return pages
    except Exception as e:
        print(f"PDF处理错误: {e}")
        return []

def extract_text_from_ppt(file_path):
    """从PPT提取文本，返回每页的文本内容列表"""
    try:
        prs = Presentation(file_path)
        slides = []
        
        for slide in prs.slides:
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            slides.append("\n".join(text_content))
            
        return slides
    except Exception as e:
        print(f"PPT处理错误: {e}")
        return []

def group_events(events, max_array_num):
    """
    根据max_array_num对events进行分组
    例如: events长度为7, max_array_num=3, 则分成[0,1,2], [3,4,5], [6]三组
    """
    if not max_array_num or max_array_num <= 1:
        return events
    
    grouped_events = []
    for i in range(0, len(events), max_array_num):
        group = events[i:i+max_array_num]
        if len(group) > 0:
            grouped_events.append('\n'.join(group))
    
    return grouped_events

def get_file_extension(file_path):
    """获取文件扩展名（小写）"""
    _, ext = os.path.splitext(file_path)
    return ext.lower()

# **Function definition**
def run_node(node):
    Array = []
    events = []
    
    # 读取输入文件路径和参数
    file_path = node['Inputs'][0]['Context']
    divide_sign = node['Inputs'][1]['Context']
    max_array_num_input = node['Inputs'][2]['Context']
    
    # 转换max_array_num为整数，默认为1
    try:
        max_array_num = int(max_array_num_input) if max_array_num_input else 1
    except (ValueError, TypeError):
        max_array_num = 1
    
    print("文件路径:", file_path)  # 调试信息，查看文件路径
    
    try:
        file_extension = get_file_extension(file_path)
        
        # 根据文件类型采用不同的处理方式
        if file_extension == '.pdf':
            # PDF文件处理
            print("正在处理PDF文件...")
            if divide_sign and divide_sign.lower() == 'page':
                # 以页为单位分割
                events = extract_text_from_pdf(file_path)
            else:
                # 将整个PDF内容合并然后按指定分隔符分割
                all_text = '\n\n'.join(extract_text_from_pdf(file_path))
                if not divide_sign or divide_sign.lower() in ['none', 'null', '*n', '']:
                    divide_sign = '\n\n'  # 默认分隔符为双回车
                
                # 检查是否有自定义的分隔符格式
                custom_pattern = parse_custom_pattern(divide_sign) if divide_sign else None
                
                if custom_pattern:
                    events = split_by_custom_pattern(all_text, custom_pattern)
                else:
                    events = all_text.split(divide_sign)
                
        elif file_extension in ['.ppt', '.pptx']:
            # PPT文件处理
            print("正在处理PPT文件...")
            if divide_sign and divide_sign.lower() == 'page':
                # 以幻灯片为单位分割
                events = extract_text_from_ppt(file_path)
            else:
                # 将整个PPT内容合并然后按指定分隔符分割
                all_text = '\n\n'.join(extract_text_from_ppt(file_path))
                if not divide_sign or divide_sign.lower() in ['none', 'null', '*n', '']:
                    divide_sign = '\n\n'  # 默认分隔符为双回车
                
                # 检查是否有自定义的分隔符格式
                custom_pattern = parse_custom_pattern(divide_sign) if divide_sign else None
                
                if custom_pattern:
                    events = split_by_custom_pattern(all_text, custom_pattern)
                else:
                    events = all_text.split(divide_sign)
                
        else:
            # 文本文件处理 (原有逻辑)
            print("正在处理文本文件...")
            # 尝试多种编码方式读取文件
            encodings_to_try = ['utf-8', 'gb2312', 'gbk', 'gb18030', 'ascii']
            content = None
            
            for encoding in encodings_to_try:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    print(f"成功使用 {encoding} 编码读取文件")
                    break  # 如果成功读取，跳出循环
                except UnicodeDecodeError:
                    print(f"使用 {encoding} 编码读取失败，尝试下一个编码")
                    continue
            
            if content is None:
                raise Exception("无法使用任何已知编码读取文件")
            
            # 检查是否以页为单位分割
            if divide_sign and divide_sign.lower() == 'page':
                # 文本文件没有明确的页概念，使用双回车作为默认分隔符
                print("文本文件不支持按页分割，使用双回车分割")
                events = content.split('\n\n')
            else:
                # 使用原有的分割逻辑
                # 检查是否有自定义的分隔符格式
                custom_pattern = parse_custom_pattern(divide_sign) if divide_sign else None
                
                if custom_pattern:
                    events = split_by_custom_pattern(content, custom_pattern)
                else:
                    # 使用原有的分割逻辑
                    if not divide_sign or divide_sign.lower() in ['none', 'null', '*n', '']:
                        divide_sign = '\n\n'  # 默认分隔符为双回车
                    events = content.split(divide_sign)
        
        # 过滤掉空白或只包含空格的事件
        events = [event.strip() for event in events if event.strip()]
        
        if not events:
            print("Warning: No valid content found in file after filtering")
            return None
        
        # 根据max_array_num对events进行分组
        if max_array_num > 1:
            grouped_events = group_events(events, max_array_num)
        else:
            grouped_events = events
        
        # 存储事件到输出节点
        for i, event in enumerate(grouped_events):
            try:
                # 创建新的输出对象
                output_copy = copy.deepcopy(Outputs)
                
                # 更新上下文
                output_copy[0]['Context'] = event
                # 这个数组的序号
                output_copy[1]['Num'] = i  # 记录当前事件在数组中的位置
                # 这个数组的总长度
                output_copy[3]['Num'] = len(grouped_events)  # 记录数组的总长度
                
                # 如果是最后一个事件
                if i == len(grouped_events) - 1:
                    output_copy[2]['Boolean'] = True
                
                # 添加独立的输出对象
                Array.append(output_copy)
            except Exception as e:
                print(f"Error processing event: {e}")

    except FileNotFoundError:
        print(f"File not found: {file_path}")
        return None
    except Exception as e:
        print(f"An error occurred: {e}")
        return None

    # 返回 Array
    print('分割后的事件数组:', Array)
    return Array
